import os
import logging
from typing import List, Dict
import pdfplumber
import docx2txt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from config import DATA_DIR

logging.getLogger("pdfminer").setLevel(logging.ERROR)

class DocumentLoader:
    def __init__(self, data_dir: str = DATA_DIR):
        self.data_dir = data_dir
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt"]

    def _extract_shape_text_recursive(self, shape) -> str:
        """递归提取 PPT 形状中的文本，支持组合图"""
        text_parts = []

        # 1. 提取文本框内容
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                p_text = "".join(run.text for run in paragraph.runs)
                if p_text.strip():
                    text_parts.append(p_text)

        # 2. 提取表格内容
        if shape.has_table:
            for row in shape.table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text_frame:
                        cell_text = cell.text_frame.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                if row_text:
                    text_parts.append(" | ".join(row_text))

        # 3. 递归处理组合图形 (Group)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child_shape in shape.shapes:
                child_text = self._extract_shape_text_recursive(child_shape)
                if child_text:
                    text_parts.append(child_text)

        return "\n".join(text_parts)

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件"""
        results = []
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                # 遍历幻灯片中的所有形状
                for shape in slide.shapes:
                    text = self._extract_shape_text_recursive(shape)
                    if text.strip():
                        slide_texts.append(text)
                
                # 提取备注页面内容
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        slide_texts.append(f"\n[备注: {notes}]")

                page_content = "\n".join(slide_texts)
                # 即使页面没有文字，也生成一个标记，保证页码对应
                formatted_text = f"--- 幻灯片 {i + 1} ---\n{page_content}\n" if page_content else f"--- 幻灯片 {i + 1} ---\n(无文字)\n"
                results.append({"text": formatted_text})
        except Exception as e:
            print(f"读取PPTX文件出错 {file_path}: {e}")
        return results

    def load_pdf(self, file_path: str) -> List[Dict]:
        """加载PDF文件 (仅提取可选文本)"""
        results = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    formatted_text = f"--- 第 {i + 1} 页 ---\n{text}\n"
                    results.append({"text": formatted_text})
        except Exception as e:
            print(f"读取PDF文件出错 {file_path}: {e}")
        return results

    def load_docx(self, file_path: str) -> str:
        try:
            return docx2txt.process(file_path)
        except Exception as e:
            print(f"读取DOCX文件出错 {file_path}: {e}"); return ""

    def load_txt(self, file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f: return f.read()
        except Exception:
            try:
                with open(file_path, 'r', encoding='gbk') as f: return f.read()
            except Exception: return ""

    def load_document(self, file_path: str) -> List[Dict[str, str]]:
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []

        if ext == ".pdf":
            pages = self.load_pdf(file_path)
            for page_idx, page_data in enumerate(pages, 1):
                documents.append({"content": page_data["text"], "filename": filename, "filepath": file_path, "filetype": ext, "page_number": page_idx})
        elif ext == ".pptx":
            slides = self.load_pptx(file_path)
            for slide_idx, slide_data in enumerate(slides, 1):
                documents.append({"content": slide_data["text"], "filename": filename, "filepath": file_path, "filetype": ext, "page_number": slide_idx})
        elif ext == ".docx":
            content = self.load_docx(file_path)
            if content: documents.append({"content": content, "filename": filename, "filepath": file_path, "filetype": ext, "page_number": 0})
        elif ext == ".txt":
            content = self.load_txt(file_path)
            if content: documents.append({"content": content, "filename": filename, "filepath": file_path, "filetype": ext, "page_number": 0})
        return documents

    def load_all_documents(self) -> List[Dict[str, str]]:
        if not os.path.exists(self.data_dir):
            print(f"数据目录不存在: {self.data_dir}")
            return None
        documents = []
        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                if file.startswith("~$") or file.startswith("."): continue # 忽略临时文件
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    print(f"正在加载: {file_path}")
                    doc_chunks = self.load_document(file_path)
                    if doc_chunks: documents.extend(doc_chunks)
        return documents
